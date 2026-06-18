import logging
from flask import Blueprint, send_file, jsonify, current_app, request
from app import db
from app.models import Project
import os

logger = logging.getLogger(__name__)
ppt_bp = Blueprint("ppt", __name__)



@ppt_bp.route("/api/projects/<int:project_id>/scan-ppt-placeholders", methods=["GET"])
def scan_ppt_placeholders(project_id):
    """Scan uploaded PPT template and return found placeholders with actual values."""
    from app.models import Upload
    from app.services.ppt_generator import _build_replacements
    import re
    from pptx import Presentation
    
    project = Project.query.get_or_404(project_id)
    uploaded = Upload.query.filter_by(project_id=project_id, file_type="ppt_template").first()
    if not uploaded or not os.path.exists(uploaded.stored_path):
        return jsonify({"error": "No PPT template uploaded. Upload a template first."}), 400
    
    prs = Presentation(uploaded.stored_path)
    placeholders = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    found = re.findall(r'\{\{[^}]+\}\}', para.text)
                    for ph in found: placeholders.add(ph)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        found = re.findall(r'\{\{[^}]+\}\}', cell.text)
                        for ph in found: placeholders.add(ph)
    
    # Build actual values from project data
    auto_values = _build_replacements(project)
    
    known_fields = {
        "{{capacity}}": "System Capacity (kWp)", "{{flh}}": "Full Load Hours",
        "{{inv_count}}": "Inverter Count", "{{module_count}}": "Module Count",
        "{{total_investment}}": "Total Investment (PHP)", "{{year1_revenue}}": "First Year Revenue (PHP)",
        "{{payback_period}}": "Payback Period (Years)", "{{rev_20y}}": "20-Year Revenue (PHP)",
        "{{irr_20y}}": "20-Year IRR", "{{rev_5y}}": "5-Year Revenue (PHP)", "{{irr_5y}}": "5-Year IRR",
        "{{inv_power}}": "Inverter Power (kW)", "{{carbon_reduction}}": "CO2 Reduction (tons)",
        "{{pro_own}}": "Customer Name", "{{adr}}": "Address", "{{geographic}}": "Coordinates (DMS)",
        "{{ghi}}": "GHI (kWh/m²)", "{{roof_area}}": "Roof Area (m²)", "{{Exchg}}": "Exchange Rate",
        "{{Exchg_date}}": "Exchange Rate Date", "{{set}}": "Inverter Unit",
        "{{project_name}}": "Project Name", "{{project name}}": "Project Name", "{{month}}": "Month",
        "{{date}}": "Date"
    }
    
    result = []
    for ph in sorted(placeholders):
        clean = ph.replace("{", "").replace("}", "").strip()
        field = known_fields.get(ph, "")
        val = auto_values.get(ph, "")
        result.append({
            "placeholder": ph,
            "field": field,
            "auto_value": str(val) if val else "",
            "matched": bool(field) and bool(val),
            "editable": True
        })
    
    return jsonify({"placeholders": result, "count": len(result)}), 200

@ppt_bp.route("/api/projects/<int:project_id>/upload-ppt-template", methods=["POST"])
def upload_ppt_template(project_id):
    """Upload a PPT template (global - shared across all projects)."""
    from app.models import Upload
    from app import db
    
    project = Project.query.get_or_404(project_id)
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400
    file = request.files["file"]
    if not file.filename.endswith('.pptx'):
        return jsonify({"error": "Only .pptx files allowed"}), 400
    
    # Save globally (replace any existing template)
    upload_dir = os.path.join(current_app.config["UPLOAD_FOLDER"], "_global_ppt")
    os.makedirs(upload_dir, exist_ok=True)
    filepath = os.path.join(upload_dir, "template.pptx")
    file.save(filepath)
    
    # Delete old per-project records, create a global record
    Upload.query.filter_by(file_type="ppt_template").delete()
    record = Upload(project_id=project.id, file_type="ppt_template",
                    original_filename=file.filename, stored_path=filepath)
    db.session.add(record)
    db.session.commit()
    
    return jsonify({"message": "PPT template uploaded (global)", "path": filepath}), 200

@ppt_bp.route("/api/projects/<int:project_id>/generate-ppt", methods=["POST"])
def generate_ppt(project_id):
    project = Project.query.get_or_404(project_id)

    # Delegate to PPT generator (stub - Phase 5)
    from app.services.ppt_generator import generate_proposal
    output_path = generate_proposal(project, current_app.config)

    project.status = "completed"
    db.session.commit()

    return jsonify({
        "message": "PPT generated",
        "download_url": f"/api/projects/{project_id}/download-ppt",
    }), 200

@ppt_bp.route("/api/projects/<int:project_id>/download-ppt", methods=["GET"])
def download_ppt(project_id):
    project = Project.query.get_or_404(project_id)
    output_dir = current_app.config["OUTPUT_FOLDER"]
    filename = f"{project.name}_proposal.pptx".replace(" ", "_")
    filepath = os.path.join(output_dir, str(project_id), filename)
    if not os.path.exists(filepath):
        return jsonify({"error": "PPT not yet generated"}), 404
    return send_file(filepath, as_attachment=True, download_name=filename)
