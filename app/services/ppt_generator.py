import os
import shutil
import logging
from pptx import Presentation
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

PLACEHOLDER_MAP = {
    "{{capacity}}": "capacity_kwp",
    "{{flh}}": "flh",
    "{{inv_count}}": "inv_count",
    "{{module_count}}": "module_count",
    "{{total_investment}}": "total_investment",
    "{{year1_revenue}}": "year1_revenue",
    "{{payback_period}}": "payback_period",
    "{{rev_20y}}": "revenue_20y",
    "{{irr_20y}}": "irr_20y",
    "{{rev_5y}}": "revenue_5y",
    "{{irr_5y}}": "irr_5y",
    "{{inv_power}}": "inv_power",
    "{{carbon_reduction}}": "carbon_reduction",
    "{{pro_own}}": "customer_name",
    "{{adr}}": "address",
    "{{geographic}}": "geographic",
    "{{ghi}}": "ghi",
    "{{roof_area}}": "roof_area",
    "{{Exchg}}": "exchange_rate",
    "{{Exchg_date}}": "exchange_date",
    "{{project name}}": "project_name",
    "{{project_name}}": "project_name",
    "{{month}}": "month",
    "{{date}}": "date",
    "{{set}}": "set_unit",
}

def generate_proposal(project, config, custom_values=None) -> str:
    """Generate proposal PPT by replacing placeholders. (Phase 5)"""
    output_folder = config["OUTPUT_FOLDER"]
    from app.models import Upload
    output_folder = config["OUTPUT_FOLDER"]

    uploaded = Upload.query.filter_by(project_id=project.id, file_type="ppt_template").first()
    if not uploaded or not os.path.exists(uploaded.stored_path):
        raise FileNotFoundError("No PPT template uploaded. Go to Step 6 and upload a template first.")
    template_path = uploaded.stored_path
    # Replace text in all shapes (handles multi-run placeholders)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if not para.runs: continue
                    full = para.text
                    newtext = full
                    for ph, val in replacements.items():
                        if ph in newtext:
                            newtext = newtext.replace(ph, str(val))
                    if newtext != full:
                        para.runs[0].text = newtext
                        for r in para.runs[1:]: r.text = ""
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            if not para.runs: continue
                            full = para.text
                            newtext = full
                            for ph, val in replacements.items():
                                if ph in newtext:
                                    newtext = newtext.replace(ph, str(val))
                            if newtext != full:
                                para.runs[0].text = newtext
                                for r in para.runs[1:]: r.text = ""

    # Apply custom overrides
    if custom_values:
        for ph, val in custom_values.items():
            full_ph = "{{" + ph + "}}" if not ph.startswith("{{") else ph
            if val:
                replacements[full_ph] = val

    prs.save(output_path)
    logger.info("PPT saved to %s", output_path)
    return output_path

def _build_replacements(project) -> dict:
    """Build placeholder replacement dictionary from project data."""
    pv = project.pv_system
    sr = project.solar_resource
    fr = None
    if project.financial_results:
        fr = project.financial_results[0]

    # Module / inverter counts
    module_count = sum(m.quantity for m in project.modules) if project.modules else 0
    inv_count = sum(i.quantity for i in project.inverters) if project.inverters else 0
    inv_power = project.inverters[0].capacity_kw if project.inverters else 0

    # Coordinates as DMS
    geographic = ""
    if project.latitude and project.longitude:
        def to_dms(deg, is_lat):
            d = int(abs(deg))
            m = int((abs(deg) - d) * 60)
            s = (abs(deg) - d - m / 60) * 3600
            direction = "N" if is_lat and deg >= 0 else "S" if is_lat else "E" if deg >= 0 else "W"
            return f"{d}°{m}'{s:.2f}\"{direction}"
        lat_dms = to_dms(project.latitude, True)
        lng_dms = to_dms(project.longitude, False)
        geographic = f"{lat_dms} {lng_dms}"

    return {
        "{{capacity}}": f"{pv.capacity_kwp:.2f}" if pv and pv.capacity_kwp else "",
        "{{flh}}": f"{sr.pvspecific_kwh_kwp:.0f}" if sr and sr.pvspecific_kwh_kwp else "",
        "{{inv_count}}": str(inv_count),
        "{{module_count}}": str(module_count),
        "{{total_investment}}": f"{fr.total_investment_php:,.0f}" if fr and fr.total_investment_php else "",
        "{{year1_revenue}}": f"{fr.year1_revenue_php:,.0f}" if fr and fr.year1_revenue_php else "",
        "{{payback_period}}": f"{fr.payback_period_years:.2f}" if fr and fr.payback_period_years else "",
        "{{rev_20y}}": f"{fr.revenue_20y_php:,.0f}" if fr and fr.revenue_20y_php else "",
        "{{irr_20y}}": f"{fr.irr*100:.2f}%" if fr and fr.irr else "",
        "{{rev_5y}}": f"{fr.revenue_5y_php:,.0f}" if fr and fr.revenue_5y_php else "",
        "{{irr_5y}}": f"{fr.irr*100:.2f}%" if fr and fr.irr else "",
        "{{inv_power}}": str(inv_power),
        "{{carbon_reduction}}": f"{pv.co2_reduction_tons:.0f}" if pv and pv.co2_reduction_tons else "",
        "{{pro_own}}": project.customer_name or "",
        "{{adr}}": project.address or "",
        "{{geographic}}": geographic,
        "{{ghi}}": f"{sr.ghi_kwhm2:.0f}" if sr and sr.ghi_kwhm2 else "",
        "{{roof_area}}": "1014",  # placeholder - from user input
        "{{Exchg}}": "8.76",
        "{{Exchg_date}}": "April 2026",
        "{{set}}": "sets",
    }
